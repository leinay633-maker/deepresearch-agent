from __future__ import annotations

import html
import json
from pathlib import Path
from typing import Iterable

from docx import Document
from pptx import Presentation
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

from deepresearch_agent.schemas import StructuredReport


SUPPORTED_EXPORT_FORMATS = {"markdown", "md", "html", "json", "pdf", "docx", "pptx"}


def export_report(
    report: StructuredReport,
    output_dir: str | Path,
    formats: Iterable[str] = ("markdown", "html", "json"),
) -> dict[str, str]:
    target_dir = Path(output_dir)
    target_dir.mkdir(parents=True, exist_ok=True)
    paths: dict[str, str] = {}
    for requested in formats:
        fmt = _normalize_format(requested)
        suffix = "md" if fmt == "markdown" else fmt
        path = target_dir / f"{report.run_id}.{suffix}"
        if fmt == "markdown":
            path.write_text(report_to_markdown(report), encoding="utf-8")
        elif fmt == "html":
            path.write_text(report_to_html(report), encoding="utf-8")
        elif fmt == "json":
            path.write_text(
                json.dumps(report.model_dump(mode="json"), ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        elif fmt == "pdf":
            write_report_pdf(report, path)
        elif fmt == "docx":
            write_report_docx(report, path)
        elif fmt == "pptx":
            write_report_pptx(report, path)
        paths[fmt] = str(path)
    return paths


def report_to_markdown(report: StructuredReport) -> str:
    lines = [
        f"# DeepResearch Report: {report.query}",
        "",
        f"- Run ID: `{report.run_id}`",
        f"- Citation retention: `{report.citation_check.retention_rate}`",
        f"- Total tokens: `{report.cost.total_tokens}`",
        f"- Estimated cost USD: `{report.cost.total_estimated_cost_usd}`",
        "",
        "## Answer",
        "",
        report.answer.strip(),
        "",
        "## Sources",
    ]
    for source in report.sources:
        lines.append(f"- [{source.id}] {source.title} - {source.url}")
    lines.extend(["", "## Citation Assessments"])
    for assessment in report.citation_check.assessments:
        lines.append(
            f"- `{assessment.support_level}` score `{assessment.overlap_score}`: "
            f"{assessment.claim}"
        )
        for quote in assessment.evidence_quotes:
            lines.append(f"  - Evidence [{quote.source_id}]: {quote.quote}")
    return "\n".join(lines).rstrip() + "\n"


def report_to_html(report: StructuredReport) -> str:
    sources = "\n".join(
        f"<li><strong>{html.escape(source.id)}</strong> "
        f"{html.escape(source.title)} - "
        f"<a href=\"{html.escape(source.url)}\">{html.escape(source.url)}</a></li>"
        for source in report.sources
    )
    assessments = []
    for assessment in report.citation_check.assessments:
        quotes = "".join(
            f"<li>Evidence {html.escape(quote.source_id)}: {html.escape(quote.quote)}</li>"
            for quote in assessment.evidence_quotes
        )
        assessments.append(
            "<li>"
            f"<span>{html.escape(assessment.support_level)}</span> "
            f"<code>{assessment.overlap_score}</code> "
            f"{html.escape(assessment.claim)}"
            f"<ul>{quotes}</ul>"
            "</li>"
        )
    return "\n".join(
        [
            "<!doctype html>",
            "<html lang=\"en\">",
            "<head>",
            "  <meta charset=\"utf-8\">",
            f"  <title>{html.escape(report.query)}</title>",
            "  <style>body{font-family:Segoe UI,Arial,sans-serif;max-width:960px;margin:32px auto;line-height:1.55;color:#1f2933}pre{white-space:pre-wrap;background:#f4f6f8;padding:16px;border-radius:6px}code{background:#eef1f5;padding:1px 4px;border-radius:4px}</style>",
            "</head>",
            "<body>",
            f"  <h1>{html.escape(report.query)}</h1>",
            f"  <p><strong>Run ID:</strong> <code>{html.escape(report.run_id)}</code></p>",
            f"  <p><strong>Citation retention:</strong> {report.citation_check.retention_rate}</p>",
            f"  <p><strong>Total tokens:</strong> {report.cost.total_tokens}</p>",
            f"  <p><strong>Estimated cost USD:</strong> {report.cost.total_estimated_cost_usd}</p>",
            "  <h2>Answer</h2>",
            f"  <pre>{html.escape(report.answer.strip())}</pre>",
            "  <h2>Sources</h2>",
            f"  <ul>{sources}</ul>",
            "  <h2>Citation Assessments</h2>",
            f"  <ul>{''.join(assessments)}</ul>",
            "</body>",
            "</html>",
            "",
        ]
    )


def write_report_docx(report: StructuredReport, path: Path) -> None:
    document = Document()
    document.add_heading(f"DeepResearch Report: {report.query}", level=1)
    _docx_paragraph(document, f"Run ID: {report.run_id}")
    _docx_paragraph(document, f"Citation retention: {report.citation_check.retention_rate}")
    _docx_paragraph(document, f"Total tokens: {report.cost.total_tokens}")
    _docx_paragraph(
        document,
        f"Estimated cost USD: {report.cost.total_estimated_cost_usd}",
    )
    document.add_heading("Answer", level=2)
    for paragraph in _text_blocks(report.answer):
        _docx_paragraph(document, paragraph)
    document.add_heading("Sources", level=2)
    for source in report.sources:
        _docx_paragraph(document, f"[{source.id}] {source.title} - {source.url}", style="List Bullet")
    document.add_heading("Citation Assessments", level=2)
    for assessment in report.citation_check.assessments:
        _docx_paragraph(
            document,
            f"{assessment.support_level} score {assessment.overlap_score}: {assessment.claim}",
            style="List Bullet",
        )
        for quote in assessment.evidence_quotes:
            _docx_paragraph(
                document,
                f"Evidence [{quote.source_id}]: {quote.quote}",
                style="List Bullet 2",
            )
    document.save(path)


def write_report_pdf(report: StructuredReport, path: Path) -> None:
    styles = getSampleStyleSheet()
    story = [
        Paragraph(html.escape(f"DeepResearch Report: {report.query}"), styles["Title"]),
        Spacer(1, 12),
        Paragraph(html.escape(f"Run ID: {report.run_id}"), styles["Normal"]),
        Paragraph(
            html.escape(f"Citation retention: {report.citation_check.retention_rate}"),
            styles["Normal"],
        ),
        Paragraph(html.escape(f"Total tokens: {report.cost.total_tokens}"), styles["Normal"]),
        Paragraph(
            html.escape(f"Estimated cost USD: {report.cost.total_estimated_cost_usd}"),
            styles["Normal"],
        ),
        Spacer(1, 12),
        Paragraph("Answer", styles["Heading2"]),
    ]
    for paragraph in _text_blocks(report.answer):
        story.append(Paragraph(html.escape(paragraph), styles["BodyText"]))
        story.append(Spacer(1, 6))
    story.extend([Spacer(1, 8), Paragraph("Sources", styles["Heading2"])])
    for source in report.sources:
        story.append(Paragraph(html.escape(f"[{source.id}] {source.title} - {source.url}"), styles["BodyText"]))
    story.extend([Spacer(1, 8), Paragraph("Citation Assessments", styles["Heading2"])])
    for assessment in report.citation_check.assessments:
        story.append(
            Paragraph(
                html.escape(
                    f"{assessment.support_level} score {assessment.overlap_score}: "
                    f"{assessment.claim}"
                ),
                styles["BodyText"],
            )
        )
        for quote in assessment.evidence_quotes:
            story.append(
                Paragraph(
                    html.escape(f"Evidence [{quote.source_id}]: {quote.quote}"),
                    styles["BodyText"],
                )
            )
    SimpleDocTemplate(str(path)).build(story)


def write_report_pptx(report: StructuredReport, path: Path) -> None:
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = f"DeepResearch Report\n{report.query}"
    title_slide.placeholders[1].text = (
        f"Run ID: {report.run_id}\n"
        f"Citation retention: {report.citation_check.retention_rate}\n"
        f"Total tokens: {report.cost.total_tokens}\n"
        f"Estimated cost USD: {report.cost.total_estimated_cost_usd}"
    )

    _add_pptx_bullets(
        presentation,
        "Answer",
        _text_blocks(report.answer[:1600]) or [report.answer[:1600]],
    )
    _add_pptx_bullets(
        presentation,
        "Sources",
        [f"[{source.id}] {source.title} - {source.url}" for source in report.sources[:10]],
    )
    assessment_lines: list[str] = []
    for assessment in report.citation_check.assessments[:8]:
        assessment_lines.append(
            f"{assessment.support_level} {assessment.overlap_score}: {assessment.claim}"
        )
        for quote in assessment.evidence_quotes[:2]:
            assessment_lines.append(f"Evidence [{quote.source_id}]: {quote.quote}")
    _add_pptx_bullets(presentation, "Citation Assessments", assessment_lines)
    presentation.save(path)


def _normalize_format(value: str) -> str:
    fmt = value.strip().lower()
    if fmt == "md":
        fmt = "markdown"
    if fmt not in SUPPORTED_EXPORT_FORMATS:
        expected = ", ".join(sorted(SUPPORTED_EXPORT_FORMATS))
        raise ValueError(f"unsupported export format: {value}; expected one of {expected}")
    return fmt


def _text_blocks(text: str) -> list[str]:
    return [block.strip() for block in text.strip().split("\n\n") if block.strip()]


def _docx_paragraph(document: Document, text: str, style: str | None = None) -> None:
    if style:
        document.add_paragraph(text, style=style)
    else:
        document.add_paragraph(text)


def _add_pptx_bullets(
    presentation: Presentation,
    title: str,
    items: list[str],
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for index, item in enumerate(items or ["No content."]):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = item[:700]
        paragraph.level = 0
